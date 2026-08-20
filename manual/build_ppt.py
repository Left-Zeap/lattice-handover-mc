# -*- coding: utf-8 -*-
"""生成汇报 PPT：manual/连续装载双光晶格输运模拟_汇报.pptx

依据 manual/OUTLINE.md §6 的 18 页结构、§2 数字口径、§6 配色。
用法：.venv/Scripts/python.exe manual/build_ppt.py
流程：1) matplotlib mathtext 渲染 F1-F12 公式 PNG -> manual/figures/formula_*.png
      2) python-pptx 组装 18 页 16:9 PPT
      3) 重新打开 PPT 做边界/重叠/标题验证，并用 PIL 抽查公式 PNG
"""
import io
import os
import sys

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "manual", "figures")
OUT = os.path.join(ROOT, "output")
PPTX_PATH = os.path.join(ROOT, "manual", "连续装载双光晶格输运模拟_汇报.pptx")

BLUE = RGBColor(0x1F, 0x3B, 0x5C)
RED = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_HEX = "#1F3B5C"

PAGE_W = Inches(13.333)
PAGE_H = Inches(7.5)
BAND_H = Inches(0.85)
FONT = "Microsoft YaHei"

# ---------------------------------------------------------------- 公式渲染

FORMULAS = {
    "F1": [r"$U(\mathbf{r})=\frac{\pi c^{2}}{2}\sum_{j}\frac{g_{j}\Gamma_{j}"
           r"}{\omega_{j}^{3}}\left[\frac{1}{\Delta_{j}}-\frac{1}{\omega_{L}+\omega_{j}}\right]I(\mathbf{r})$"],
    "F2": [r"$I_{\mathrm{anti}}=\frac{2P_{f}(1+\sqrt{R})^{2}}{\pi w^{2}}$",
           r"$U_{0}=A(\delta)\,P_{\mathrm{src}},\quad \Gamma_{\mathrm{sc}}=S(\delta)\,P_{\mathrm{src}}$"],
    "F3": [r"$\omega_{z}=k\sqrt{2U/m},\quad \omega_{r}=\sqrt{4U/(m w^{2})},\quad a_{c}=Uk/m$"],
    "F4": [r"$v=\lambda\,\Delta\nu\,/\,2$"],
    "F5": [r"$U_{\mathrm{eff}}/U_{0}=\sqrt{1-\beta^{2}}-\beta\,\arccos\beta,\quad \beta=a/a_{c}$"],
    "F6": [r"$F_{3}(\eta)=1-e^{-\eta}\left(1+\eta+\eta^{2}/2\right),\quad \eta=U/(k_{B}T)$"],
    "F7": [r"$T\propto\bar{\omega}\ (\mathrm{adiabatic}),\quad "
           r"\Delta T_{\mathrm{rec}}=\frac{2E_{r}}{3k_{B}}\,N_{\mathrm{sc}}$",
           r"$\Delta T_{\Delta a}=\frac{m(\Delta a)^{2}}{6k_{B}\omega_{z}^{2}},\quad "
           r"\Gamma_{\mathrm{para}}=\bar{\omega}^{2}S_{\epsilon}(2\bar{\omega})/4$"],
    "F8": [r"$m\,dv=-\nabla U_{L1}\,dt-m\gamma v\,dt+\sqrt{2m\gamma k_{B}T_{\mathrm{eq}}}\,dW+dp_{\mathrm{rec}}$",
           r"$T(t)=T_{\mathrm{eq}}+\left[T(0)-T_{\mathrm{eq}}\right]e^{-2\gamma t}$"],
    "F9": [r"$E'=\frac{1}{2}m|\mathbf{v}'|^{2}+U'(\mathbf{r}')<U_{2,\mathrm{eff}}$"],
    "F10": [r"$n_{0}=N\,\omega_{r}^{2}\,\omega_{z}\left(\frac{m}{2\pi k_{B}T}\right)^{3/2}$"],
    "F11": [r"$P\geq P_{U},\ P\geq P_{HO}\propto\tau_{HO}^{-2},\ P\geq P_{\mathrm{bound}},"
            r"\ P\leq P_{\mathrm{sc}},\ P\leq P_{\max}$",
            r"$J=P/P_{\max}+0.05\,(\delta-\delta_{\min})/(\delta_{\max}-\delta_{\min})$"],
    "F12": [r"$S_{\mathrm{total}}=\eta_{\mathrm{load}}\,S_{L1}\,\eta_{HO}\,S_{L2},"
            r"\quad S_{L1}=S_{\mathrm{spill}}\,S_{\mathrm{rate}}$"],
}


def render_line_png(math_str, fontsize=22, color=BLUE_HEX, dpi=200):
    fig = plt.figure()
    fig.text(0, 0, math_str, fontsize=fontsize, color=color)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, transparent=True,
                bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)
    buf.seek(0)
    return Image.open(buf).convert("RGBA")


def render_formulas():
    paths = {}
    for key, lines in FORMULAS.items():
        imgs = [render_line_png(s) for s in lines]
        gap = 18
        w = max(im.width for im in imgs)
        h = sum(im.height for im in imgs) + gap * (len(imgs) - 1)
        canvas = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        y = 0
        for im in imgs:
            canvas.paste(im, ((w - im.width) // 2, y), im)
            y += im.height + gap
        path = os.path.join(FIG, f"formula_{key}.png")
        canvas.save(path)
        paths[key] = path
        print(f"  公式 {key}: {canvas.size} -> {path}")
    return paths


def check_formula_pngs(paths, keys=("F1", "F5", "F9")):
    for k in keys:
        im = Image.open(paths[k]).convert("RGBA")
        alpha = im.getchannel("A")
        nonzero = sum(1 for v in alpha.getdata() if v > 10)
        frac = nonzero / (im.width * im.height)
        status = "OK" if frac > 0.01 else "EMPTY!"
        print(f"  公式 PNG 抽查 {k}: {im.size}, 非透明像素占比 {frac:.3f} -> {status}")
        assert frac > 0.01, f"formula_{k}.png 渲染为空"


# ---------------------------------------------------------------- PPT 助手

def set_text(tf, lines, size=18, color=BLUE, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.1):
    """lines: list of str 或 list of (str, dict) 覆盖样式"""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, kw = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = kw.get("align", align)
        p.line_spacing = kw.get("line_spacing", line_spacing)
        p.space_after = Pt(kw.get("space_after", 4))
        run = p.add_run()
        run.text = text
        f = run.font
        f.size = Pt(kw.get("size", size))
        f.bold = kw.get("bold", bold)
        f.color.rgb = kw.get("color", color)
        f.name = FONT
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", FONT)


def add_textbox(slide, x, y, w, h, lines, **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(box.text_frame, lines, **kw)
    return box


def add_band_and_title(slide, title):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PAGE_W, BAND_H)
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    band.shadow.inherit = False
    add_textbox(slide, 0.4, 0.08, 12.5, 0.7, [title], size=30, color=WHITE,
                bold=True, anchor=MSO_ANCHOR.MIDDLE)


def add_page_number(slide, n, total=18):
    add_textbox(slide, 12.35, 7.12, 0.85, 0.3, [f"{n} / {total}"], size=12,
                color=GRAY, align=PP_ALIGN.RIGHT)


def add_picture_fit(slide, path, x, y, w, h):
    """在目标矩形 (x,y,w,h, 单位英寸) 内保持宽高比居中贴图"""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(path, Inches(px), Inches(py),
                                    Inches(pw), Inches(ph))


def add_box(slide, x, y, w, h, lines, fill=None, line_color=None, **kw):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    set_text(shape.text_frame, lines, anchor=MSO_ANCHOR.MIDDLE, **kw)
    return shape


# ---------------------------------------------------------------- 18 页内容

def build_deck(fp):
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H
    blank = prs.slide_layouts[6]
    titles = []

    def new_slide(title=None):
        slide = prs.slides.add_slide(blank)
        if title:
            add_band_and_title(slide, title)
            titles.append(title)
        return slide

    # ---- 1 封面
    s = new_slide()
    titles.append("连续装载双光晶格输运模拟")
    add_textbox(s, 0.8, 1.9, 11.73, 1.4, ["连续装载双光晶格输运模拟"],
                size=44, color=BLUE, bold=True, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 0.8, 3.45, 11.73, 0.7, ["物理模型 · 计算方法 · 程序实现"],
                size=24, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.8, 4.35, 11.73, 0.5, ["2026-08"], size=18, color=GRAY,
                align=PP_ALIGN.CENTER)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.15), PAGE_W, Inches(1.35))
    band.fill.solid(); band.fill.fore_color.rgb = BLUE
    band.line.fill.background(); band.shadow.inherit = False
    set_text(band.text_frame,
             ["150 ms/团 交付    |    2.5×10⁶ 原子/团    |    温度 ~120 µK    |    双晶格效率 ~60%"],
             size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # ---- 2 工作目标
    s = new_slide("工作目标：面向连续 3000 量子比特的输运链路")
    tasks = ["① 复现 Rb-87 论文\n输运链路全流程", "② 设计 Cs-133 方案\n阱深 / 失谐 / 功率",
             "③ 建立全链路模拟\n与参数设计工具"]
    for i, t in enumerate(tasks):
        add_box(s, 0.55 + i * 4.15, 1.35, 3.85, 1.75, t.split("\n"),
                fill=RGBColor(0xEA, 0xEF, 0xF5), size=20, color=BLUE, bold=True,
                align=PP_ALIGN.CENTER)
    nums = [("150 ms/团", "交付节奏"), ("2.5×10⁶", "原子/团"),
            ("~120 µK", "交付温度"), ("~60%", "双晶格效率")]
    for i, (num, lab) in enumerate(nums):
        add_box(s, 0.55 + i * 3.11, 3.45, 2.87, 1.7,
                [(num, {"size": 30, "bold": True, "color": RED}),
                 (lab, {"size": 16, "color": GRAY})],
                fill=None, line_color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.55, 5.5, 12.3, 1.4,
                ["• 基准：Nature 646, 1075 (2025) —— 连续运行 3000 量子比特，科学区 >2 h 维持",
                 "• 核心问题：原子数留存 × 温度控制 × 散射/相干性 → 需要全链路数值模拟 + 参数设计"],
                size=18)

    # ---- 3 总体架构
    s = new_slide("总体架构：MOT → L1 → 交接 → L2 → 科学区")
    add_picture_fit(s, os.path.join(FIG, "arch_schematic.png"), 0.3, 1.5, 9.3, 4.6)
    add_textbox(s, 9.8, 1.7, 3.3, 4.4,
                ["• L1：39 cm / 50 ms", "• 交接：1 ms，无冷却", "• L2：17 cm / 21 ms",
                 "• DPT 倾角 ~4°", "• DPT 截断 e⁻¹⁸", "• 全程阱深 >500 µK"],
                size=18, line_spacing=1.4)
    add_page_number(s, 3)

    # ---- 4 时序总览
    s = new_slide("时序总览：一个装载周期")
    add_picture_fit(s, os.path.join(FIG, "timing_sequence.png"), 0.5, 1.2, 12.3, 5.25)
    add_textbox(s, 0.6, 6.55, 12.1, 0.45,
                ["总周期 171 ms；稳态交付节奏 150 ms/团（MOT 80 → 压缩 7 → idle 1 → LGM 11 → L1 50 → HO 1 → L2 21 ms）"],
                size=18, align=PP_ALIGN.CENTER)
    add_page_number(s, 4)

    # ---- 5 物理基础①
    s = new_slide("物理基础①：偶极势与散射率（F1 · F2）")
    add_picture_fit(s, fp["F1"], 0.25, 1.3, 5.3, 2.5)
    add_picture_fit(s, fp["F2"], 0.25, 4.0, 5.3, 2.0)
    add_picture_fit(s, os.path.join(FIG, "dipole_curves.png"), 5.75, 1.3, 7.3, 4.9)
    add_textbox(s, 0.6, 6.4, 12.1, 0.5,
                ["D1 红失谐：阱深 ∝ P/δ，散射 ∝ P/δ² —— 加大失谐以功率换相干性；D1:D2 权重 1:2"],
                size=18, align=PP_ALIGN.CENTER)
    add_page_number(s, 5)

    # ---- 6 物理基础②
    s = new_slide("物理基础②：传送带、加速势垒与热束缚（F3–F6）")
    add_picture_fit(s, fp["F4"], 0.3, 1.2, 2.7, 1.15)
    add_picture_fit(s, fp["F3"], 3.15, 1.2, 4.5, 1.15)
    add_picture_fit(s, fp["F5"], 7.8, 1.2, 5.2, 1.15)
    add_picture_fit(s, os.path.join(FIG, "conveyor_principle.png"), 0.35, 2.55, 6.3, 3.6)
    add_picture_fit(s, os.path.join(FIG, "tilted_barrier.png"), 6.85, 2.55, 6.15, 3.6)
    add_textbox(s, 0.6, 6.35, 12.1, 0.6,
                ["工作点 a ≈ 4000 m/s² ≪ a_c，倾斜势垒仍高；热束缚 η=U/(k_BT)=5 → 87.5%，η=10 → 99.95%（F6）"],
                size=18, align=PP_ALIGN.CENTER)
    add_page_number(s, 6)

    # ---- 7 计算方法总览
    s = new_slide("计算方法总览：双轨制 + LP 初筛")
    add_picture_fit(s, os.path.join(FIG, "pipeline_flowchart.png"), 1.6, 1.2, 10.1, 5.2)
    add_textbox(s, 0.6, 6.5, 12.1, 0.45,
                ["解析宏观腿（快，用于扫描） + 轨迹 Monte Carlo（准，用于核验）；失谐–功率 LP 可行域先行初筛"],
                size=18, align=PP_ALIGN.CENTER)
    add_page_number(s, 7)

    # ---- 8 装载阶段
    s = new_slide("装载：LGM 冷却 + 静止晶格（F8）")
    add_picture_fit(s, fp["F8"], 0.3, 1.4, 5.4, 2.6)
    add_textbox(s, 0.35, 4.25, 5.4, 2.0,
                ["• Langevin 动力学 + 反冲项", "• 4×10⁶ 原子 @ ~20 µK 入阱",
                 "• LGM 11 ms 装载，扫描定功率"], size=18, line_spacing=1.4)
    add_picture_fit(s, os.path.join(OUT, "loading_ramp_scan_rb87.png"), 5.95, 1.3, 7.05, 4.95)
    add_page_number(s, 8)

    # ---- 9 L1 运输
    s = new_slide("L1 运输：分项温升预算（F7）")
    add_picture_fit(s, fp["F7"], 0.3, 1.35, 5.7, 2.6)
    add_textbox(s, 0.35, 4.15, 5.55, 2.3,
                ["• 39 cm / 50 ms；v = 8–10 m/s", "• a ≈ 4000 m/s²，恒阱深功率跟随",
                 "• 路径：20 → +10.8 → +76.2* → +13.0 ≈ 120 µK", "  （*交接随机相位解析上界）"],
                size=18, line_spacing=1.35)
    add_picture_fit(s, os.path.join(OUT, "figures", "fig2_temperature_path.png"),
                    6.1, 1.35, 7.0, 4.85)
    add_page_number(s, 9)

    # ---- 10 交接①物理
    s = new_slide("交接①：双晶格时变势与捕获判据（F9）")
    add_picture_fit(s, os.path.join(FIG, "handover_potential.png"), 0.35, 1.25, 8.5, 5.3)
    add_picture_fit(s, fp["F9"], 9.05, 1.6, 3.95, 1.6)
    add_textbox(s, 9.05, 3.5, 3.95, 3.0,
                ["• 1 ms 强度反向 ramp", "• 相对相位 0 / π →", "  势垒与能量再分配",
                 "• 逐粒子 MC 判定", "  不接受解析替代"], size=18, line_spacing=1.35)
    add_page_number(s, 10)

    # ---- 11 交接②理论修正
    s = new_slide("交接②：理论修正 —— 六大错误")
    rows = [
        ("#", "错误推导", "修正口径"),
        ("1", "条件极小值 ≠ 移动阱中心", "以移动阱中心展开"),
        ("2", "小角度 ≠ 小相位", "kθσ_y ≈ 15.2 ≫ 1"),
        ("3", "忽略周期性 → 能量无界", "ΔE_site = U·sin²(kδq) ≤ U"),
        ("4", "自由度计数混用", "统一 3D 计数"),
        ("5", "1 ms 误用突然极限", "ω∥τ ≈ 1754，绝热区"),
        ("6", "U/(k_BT)≥α 当交接率", "是势深裕量；P_bound(5)=87.5%"),
    ]
    tbl_shape = s.shapes.add_table(7, 3, Inches(0.35), Inches(1.25),
                                   Inches(8.35), Inches(5.4))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(0.55)
    tbl.columns[1].width = Inches(3.9)
    tbl.columns[2].width = Inches(3.9)
    for r, row in enumerate(rows):
        for c, txt in enumerate(row):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            set_text(cell.text_frame, [txt],
                     size=14 if r else 15, bold=(r == 0),
                     color=WHITE if r == 0 else BLUE,
                     align=PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else (
                RGBColor(0xF2, 0xF5, 0xF9) if r % 2 else WHITE)
    add_textbox(s, 8.95, 1.7, 4.05, 1.5, ["3238 W → 3.5 W"], size=36, color=RED,
                bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, 8.95, 3.5, 4.05, 2.8,
                ["修正后 Cs 交接需求：", "• 阱深 614 µK", "• D1 失谐 600–700 GHz",
                 "• 散射 ~500 s⁻¹", "• 功率仅需 3–3.5 W"], size=18, line_spacing=1.3)
    add_page_number(s, 11)

    # ---- 12 交接③MC 验证
    s = new_slide("交接③：Monte Carlo 验证")
    add_picture_fit(s, os.path.join(OUT, "handover_efficiency_map.png"),
                    0.4, 1.25, 8.3, 5.35)
    add_textbox(s, 8.95, 1.7, 4.0, 1.0, ["交接率 ≈ 1.000"], size=30, color=RED,
                bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, 8.95, 2.85, 4.0, 0.9, ["净升温 ≈ 5.8 µK"], size=24, color=BLUE,
                bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, 8.95, 4.0, 4.0, 2.4,
                ["远低于解析上界 +76.2 µK：随机相位上界过于保守，", "正确工作点下交接近乎无损。"],
                size=18)
    add_page_number(s, 12)

    # ---- 13 L2 与科学区
    s = new_slide("L2 与科学区：密度与全链留存（F10 · F12）")
    add_picture_fit(s, fp["F10"], 0.3, 1.35, 5.5, 1.6)
    add_picture_fit(s, fp["F12"], 0.3, 3.15, 5.5, 1.6)
    add_textbox(s, 0.35, 4.95, 5.5, 1.6,
                ["• L2：17 cm / 21 ms", "• P_L2,end = 0.36·P_handover",
                 "• 交付 2.5×10⁶ 原子 @ ~120 µK"], size=18, line_spacing=1.35)
    add_picture_fit(s, os.path.join(OUT, "full_chain_scan_rb87.png"), 6.0, 1.3, 7.05, 4.95)
    add_page_number(s, 13)

    # ---- 14 设计优化
    s = new_slide("设计优化：失谐–功率 LP 可行域（F11）")
    add_picture_fit(s, fp["F11"], 0.3, 1.5, 5.3, 3.0)
    add_textbox(s, 0.35, 4.7, 5.3, 1.8,
                ["• 五约束半平面 → 凸可行域", "• LP 快速初筛参数区",
                 "• 再由 handover map MC 复核"], size=18, line_spacing=1.35)
    add_picture_fit(s, os.path.join(OUT, "detuning_power_lp.png"), 5.85, 1.25, 7.15, 5.3)
    add_page_number(s, 14)

    # ---- 15 Cs-133 方案
    s = new_slide("Cs-133 方案设计")
    add_picture_fit(s, os.path.join(FIG, "fig1_scheme_comparison_fixed.png"),
                    0.3, 1.25, 6.4, 5.1)
    add_picture_fit(s, os.path.join(OUT, "full_chain_scan_cs133.png"),
                    6.85, 1.25, 6.15, 5.1)
    add_textbox(s, 0.6, 6.5, 12.1, 0.45,
                ["Cs-133 工作点：阱深 614 µK · D1 失谐 600–700 GHz · 功率 3–3.5 W · 散射 ~500 s⁻¹"],
                size=18, align=PP_ALIGN.CENTER)
    add_page_number(s, 15)

    # ---- 16 程序架构
    s = new_slide("程序架构与接口")
    add_picture_fit(s, os.path.join(FIG, "module_layers.png"), 0.4, 1.25, 7.5, 5.35)
    add_textbox(s, 8.2, 1.7, 4.8, 4.6,
                ["• 两种接口：标量 (N, T)", "  与相空间连续", "  ParticleEnsemble",
                 "• 配置流：data/*.json →", "  dataclass 默认 → CLI 覆盖",
                 "• 入口：python -m", "  continuous_loading <子命令>",
                 "• GPU 支路用于批量 MC"], size=18, line_spacing=1.3)
    add_page_number(s, 16)

    # ---- 17 口径纪律
    s = new_slide("口径纪律与已知边界")
    rules = [
        ("纪律一", "数字一律采用 OUTLINE §2 口径；冲突时标注来源（论文 / 解析 / MC）"),
        ("纪律二", "论文未公开量 = 工程假设，显式标注、可追溯"),
        ("纪律三", "损失系数默认全零；交接率只认轨迹 MC，不接受解析替代"),
    ]
    for i, (tag, txt) in enumerate(rules):
        add_box(s, 0.6, 1.45 + i * 1.7, 12.1, 1.45,
                [(tag, {"size": 20, "bold": True, "color": RED}),
                 (txt, {"size": 22, "bold": True, "color": BLUE})],
                fill=RGBColor(0xF2, 0xF5, 0xF9))
    add_textbox(s, 0.6, 6.6, 12.1, 0.45,
                ["已知边界：再热化假设；随机相位解析上界 ≠ MC 真实升温"],
                size=16, color=GRAY, align=PP_ALIGN.CENTER)
    add_page_number(s, 17)

    # ---- 18 总结
    s = new_slide("总结与下一步")
    add_box(s, 0.5, 1.35, 6.0, 5.2,
            [("成果", {"size": 24, "bold": True, "color": BLUE, "align": PP_ALIGN.CENTER}),
             ("• Rb-87 全链复现论文：150 ms/团、", {"size": 18}),
             ("  2.5×10⁶ 原子 @ ~120 µK、效率 ~60%", {"size": 18}),
             ("• 交接理论修正：3238 W → 3.5 W，", {"size": 18}),
             ("  MC 交接率 ≈ 1.000、净升温 5.8 µK", {"size": 18}),
             ("• Cs-133 方案：614 µK / 600–700 GHz", {"size": 18}),
             ("  / 3–3.5 W，散射 ~500 s⁻¹", {"size": 18})],
            fill=RGBColor(0xF2, 0xF5, 0xF9), line_spacing=1.25)
    add_box(s, 6.85, 1.35, 6.0, 5.2,
            [("下一步", {"size": 24, "bold": True, "color": RED, "align": PP_ALIGN.CENTER}),
             ("• 交接 MC 与实验参数逐点对标", {"size": 18}),
             ("• 损失系数标定（当前默认全零）", {"size": 18}),
             ("• 双轨结果对照手册补齐（口径表）", {"size": 18}),
             ("• GPU 批量扫描，接入实验闭环", {"size": 18}),
             ("• 科学区装载镊阵（30 万原子/s）", {"size": 18}),
             ("  接口参数联调", {"size": 18})],
            fill=RGBColor(0xFB, 0xF0, 0xEF), line_spacing=1.25)
    add_page_number(s, 18)

    prs.save(PPTX_PATH)
    print(f"PPT 已保存：{PPTX_PATH}（{len(prs.slides.__iter__.__self__._sldIdLst)} 页）")
    return PPTX_PATH, titles


# ---------------------------------------------------------------- 验证

def verify(path):
    prs = Presentation(path)
    pw, ph = prs.slide_width, prs.slide_height
    tol = Emu(int(0.02 * 914400))  # 0.02 in 容差
    problems = []
    for idx, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        # (1) 边界
        for sh in shapes:
            if sh.left < -tol or sh.top < -tol or \
               sh.left + sh.width > pw + tol or sh.top + sh.height > ph + tol:
                problems.append(f"第{idx}页 越界: {sh.shape_type} "
                                f"({sh.left},{sh.top},{sh.width},{sh.height})")
        # (2) 图片 vs 文本框重叠（背景色带等自选图形豁免）
        pics = [sh for sh in shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
        txts = [sh for sh in shapes
                if sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and sh.has_text_frame
                and sh.text_frame.text.strip()]
        for p in pics:
            for t in txts:
                ox = min(p.left + p.width, t.left + t.width) - max(p.left, t.left)
                oy = min(p.top + p.height, t.top + t.height) - max(p.top, t.top)
                if ox > tol and oy > tol:
                    problems.append(f"第{idx}页 图片与文本框重叠: "
                                    f"文本='{t.text_frame.text[:20]}'")
        # (3) 标题存在：标题栏色带 + 非空文本
        has_band = any(sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                       and sh.top <= tol and sh.width >= pw - tol for sh in shapes)
        has_text = any(sh.has_text_frame and sh.text_frame.text.strip()
                       for sh in shapes)
        if not (has_text and (has_band or idx == 1)):
            problems.append(f"第{idx}页 缺标题/标题栏")
        print(f"  第{idx:2d}页: {len(shapes)} shapes, "
              f"{len(pics)} 图, {len(txts)} 文本框 -> 检查完成")
    if problems:
        print("验证失败：")
        for p in problems:
            print("  -", p)
        sys.exit(1)
    print(f"验证通过：{idx} 页全部满足边界 / 无图-文重叠 / 标题齐全")


if __name__ == "__main__":
    print("1) 渲染公式 PNG ...")
    fp = render_formulas()
    check_formula_pngs(fp)
    print("2) 组装 PPT ...")
    path, _ = build_deck(fp)
    print("3) 验证 PPT ...")
    verify(path)
